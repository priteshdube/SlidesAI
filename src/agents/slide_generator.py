import os
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv


class SlideGenerator:
    def __init__(self, output_dir= "output/"):
        self.output_dir = output_dir

        os.makedirs(self.output_dir, exist_ok=True)

    def create_slide_deck(self, title, bullet_points):
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(title_slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title


        for i in range(0, len(bullet_points), 5):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape= slide.shapes.title
            content_shape= slide.placeholders[1]


            title_shape.text = f"{title} (part{i//5 + 1})"
            for bullet in bullet_points[i: i+5]:
                p = content_shape.text_frame.add_paragraph()
                p.text = bullet
                p.font.size = Pt(24)
                p.level = 0

            pptx_path = os.path.join(self.output_dir, "presentation.pptx")
            prs.save(pptx_path)
            print(f"Slide deck saved to {pptx_path}")



if __name__ == "__main__":
    generator = SlideGenerator()
    title = "AI Agent for Presentations"
    bullet_points = [
        "Builds an AI agent using the Gemini API to automatically generate professional slide decks.",
        "Parses an input document, structures the content, and creates a complete presentation file.",
        "Leverages various tools for tasks like image generation, data visualization, and final document formatting.",
        "Dramatically reduces the time and effort required to create a polished presentation."
    ]
    generator.create_slide_deck(title, bullet_points)









